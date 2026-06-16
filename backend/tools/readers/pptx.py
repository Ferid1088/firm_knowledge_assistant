"""PPTXReaderTool — extract text and structure from PowerPoint PPTX files."""
from __future__ import annotations

import logging
from datetime import datetime
from pathlib import Path

from backend.core.tool_base import ToolMetadata
from backend.tools.readers.base import FileReaderTool, RawContent

logger = logging.getLogger(__name__)


class PPTXReaderTool(FileReaderTool):
    metadata = ToolMetadata(
        name="reader:pptx",
        version="1.0.0",
        tool_type="reader",
        description="Extract text and slide content from PowerPoint PPTX files",
        dependencies=["python-pptx>=0.6"],
        capabilities={
            "formats": [".pptx", ".PPTX"],
            "max_size": 100 * 1024 * 1024,
            "supports_page_detection": True,
        },
        is_production=True,
        is_async=True,
    )

    supported_extensions = [".pptx", ".PPTX"]
    format_name = "pptx"

    def _validate_dependencies(self) -> None:
        # pip: python-pptx, import: pptx
        import importlib
        try:
            importlib.import_module("pptx")
        except ImportError:
            raise ImportError(
                f"Tool '{self.metadata.name}' requires 'python-pptx>=0.6'. "
                "Install with: pip install python-pptx"
            )

    async def execute(self, input_data: str, **kwargs) -> RawContent:
        file_path = str(input_data)
        is_valid, error = await self.validate_input(file_path)
        if not is_valid:
            return RawContent(
                format="pptx",
                content="",
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                is_corrupted=True,
                warnings=[error],
            )

        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []
            pages = []

            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_texts.append(" | ".join(cells))

                slide_text = "\n".join(slide_texts)
                pages.append({
                    "number": i + 1,
                    "text": slide_text,
                    "tables": [],
                })
                if slide_text:
                    text_parts.append(f"\n--- Slide {i + 1} ---\n{slide_text}\n")

            p = Path(file_path)
            return RawContent(
                format="pptx",
                content="".join(text_parts),
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                page_count=len(prs.slides),
                pages=pages,
                file_size=p.stat().st_size,
                modified_time=datetime.fromtimestamp(p.stat().st_mtime).isoformat(),
            )

        except Exception as e:
            return self._handle_error(e, file_path)
